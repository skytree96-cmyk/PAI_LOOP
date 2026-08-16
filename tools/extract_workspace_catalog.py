"""Build a read-only catalog of the PAI_LOOP source workspace.

The script never edits source artifacts. Extracted text is written under the
repository's ignored ``.local`` directory so confidential source material is
not committed accidentally.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree as ET


TEXT_EXTENSIONS = {".md", ".txt", ".csv", ".json", ".yml", ".yaml"}
SKIP_NAMES = {"secrets.txt", ".ds_store", "thumbs.db"}


@dataclass
class CatalogEntry:
    path: str
    extension: str
    size_bytes: int
    sha256: str
    status: str = "pending"
    extracted_chars: int = 0
    metadata: dict[str, Any] = field(default_factory=dict)
    error: str | None = None
    extracted_file: str | None = None


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def normalize_text(parts: Iterable[str]) -> str:
    text = "\n".join(part.strip() for part in parts if part and part.strip())
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    from docx import Document

    document = Document(path)
    parts: list[str] = [paragraph.text for paragraph in document.paragraphs]
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"[TABLE {table_index}]")
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return normalize_text(parts), {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "sections": len(document.sections),
    }


def extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts: list[str] = []
    page_lengths: list[int] = []
    for number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        page_lengths.append(len(text))
        parts.append(f"[PAGE {number}]\n{text}")
    return normalize_text(parts), {
        "pages": len(reader.pages),
        "page_text_chars": page_lengths,
        "encrypted": bool(reader.is_encrypted),
    }


def extract_xlsx(path: Path) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=False)
    parts: list[str] = []
    sheets: list[dict[str, Any]] = []
    for worksheet in workbook.worksheets:
        parts.append(f"[SHEET {worksheet.title}]")
        nonempty_rows = 0
        for row in worksheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                nonempty_rows += 1
                parts.append("\t".join(values))
        sheets.append(
            {
                "title": worksheet.title,
                "max_row": worksheet.max_row,
                "max_column": worksheet.max_column,
                "nonempty_rows": nonempty_rows,
            }
        )
    workbook.close()
    return normalize_text(parts), {"sheets": sheets}


def extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    slide_shapes: list[int] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[SLIDE {number}]")
        slide_shapes.append(len(slide.shapes))
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return normalize_text(parts), {
        "slides": len(presentation.slides),
        "slide_shape_counts": slide_shapes,
    }


def xml_text(root: ET.Element) -> list[str]:
    return [node.text for node in root.iter() if node.text and node.text.strip()]


def extract_hwpx(path: Path) -> tuple[str, dict[str, Any]]:
    parts: list[str] = []
    section_names: list[str] = []
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        candidates = sorted(
            name
            for name in names
            if name.lower().startswith("contents/section") and name.lower().endswith(".xml")
        )
        for name in candidates:
            section_names.append(name)
            root = ET.fromstring(archive.read(name))
            parts.append(f"[SECTION {name}]")
            parts.extend(xml_text(root))
    return normalize_text(parts), {"sections": section_names}


def inspect_zip(path: Path) -> tuple[str, dict[str, Any]]:
    with zipfile.ZipFile(path) as archive:
        files = [info for info in archive.infolist() if not info.is_dir()]
        unsafe = [
            info.filename
            for info in files
            if Path(info.filename).is_absolute() or ".." in Path(info.filename).parts
        ]
        listing = [f"{info.filename}\t{info.file_size}" for info in files]
    return normalize_text(listing), {
        "members": len(files),
        "uncompressed_bytes": sum(info.file_size for info in files),
        "unsafe_member_names": unsafe,
    }


def extract_file(path: Path) -> tuple[str, dict[str, Any], str]:
    extension = path.suffix.lower()
    if extension in TEXT_EXTENSIONS:
        return decode_text(path.read_bytes()), {}, "extracted"
    if extension == ".docx":
        text, metadata = extract_docx(path)
        return text, metadata, "extracted"
    if extension in {".xlsx", ".xlsm"}:
        text, metadata = extract_xlsx(path)
        return text, metadata, "extracted"
    if extension == ".pdf":
        text, metadata = extract_pdf(path)
        return text, metadata, "extracted"
    if extension == ".pptx":
        text, metadata = extract_pptx(path)
        return text, metadata, "extracted"
    if extension == ".hwpx":
        text, metadata = extract_hwpx(path)
        return text, metadata, "extracted"
    if extension == ".zip":
        text, metadata = inspect_zip(path)
        return text, metadata, "inspected"
    if extension == ".hwp":
        return "", {"reason": "legacy HWP5 binary; use paired PDF or Hancom conversion worker"}, "deferred"
    return "", {"reason": "unsupported binary type"}, "skipped"


def iter_source_files(source_root: Path, repository_root: Path) -> Iterable[Path]:
    for path in sorted(source_root.rglob("*")):
        if not path.is_file():
            continue
        if repository_root == path or repository_root in path.parents:
            continue
        if path.name.lower() in SKIP_NAMES:
            continue
        if any(part in {".git", ".local", "__pycache__"} for part in path.parts):
            continue
        yield path


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-root", type=Path, required=True)
    parser.add_argument("--repository-root", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    arguments = parser.parse_args()

    source_root = arguments.source_root.resolve()
    repository_root = arguments.repository_root.resolve()
    output_root = arguments.output.resolve()
    extracted_root = output_root / "extracted"
    extracted_root.mkdir(parents=True, exist_ok=True)

    entries: list[CatalogEntry] = []
    for path in iter_source_files(source_root, repository_root):
        relative = path.relative_to(source_root).as_posix()
        digest = sha256_file(path)
        entry = CatalogEntry(
            path=relative,
            extension=path.suffix.lower() or "[none]",
            size_bytes=path.stat().st_size,
            sha256=digest,
        )
        try:
            text, metadata, status = extract_file(path)
            entry.status = status
            entry.metadata = metadata
            entry.extracted_chars = len(text)
            if text:
                extracted_file = extracted_root / f"{digest}.txt"
                extracted_file.write_text(text, encoding="utf-8")
                entry.extracted_file = extracted_file.relative_to(output_root).as_posix()
        except Exception as exc:  # catalog every failure instead of aborting the run
            entry.status = "error"
            entry.error = f"{type(exc).__name__}: {exc}"
        entries.append(entry)
        print(f"{entry.status:9} {relative}")

    by_extension: dict[str, int] = {}
    by_status: dict[str, int] = {}
    for entry in entries:
        by_extension[entry.extension] = by_extension.get(entry.extension, 0) + 1
        by_status[entry.status] = by_status.get(entry.status, 0) + 1

    payload = {
        "source_root": str(source_root),
        "repository_root": str(repository_root),
        "file_count": len(entries),
        "by_extension": dict(sorted(by_extension.items())),
        "by_status": dict(sorted(by_status.items())),
        "entries": [asdict(entry) for entry in entries],
    }
    output_root.mkdir(parents=True, exist_ok=True)
    catalog_path = output_root / "workspace_catalog.json"
    catalog_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({key: payload[key] for key in ("file_count", "by_extension", "by_status")}, ensure_ascii=False))
    print(f"catalog={catalog_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
